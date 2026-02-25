
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from reportlab.lib.pagesizes import letter, landscape
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet

def generate_pdf(csv_file, output_file):
    print(f"Generating PDF from {csv_file}...")
    df = pd.read_csv(csv_file)
    data = [df.columns.tolist()] + df.values.tolist()
    
    doc = SimpleDocTemplate(output_file, pagesize=landscape(letter))
    elements = []
    
    styles = getSampleStyleSheet()
    elements.append(Paragraph("Logical Reasoning Robustness Report", styles['Title']))
    
    # Calculate column widths to avoid overflow
    # Landscape letter is 11 inches wide. Let's use 10 inches for the table.
    num_cols = len(df.columns)
    col_widths = [10*72/num_cols] * num_cols
    
    t = Table(data, colWidths=col_widths)
    t.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 8), # Small font to prevent overlap
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    elements.append(t)
    doc.build(elements)
    print(f"PDF saved to {output_file}")

def generate_pptx(csv_file, output_file):
    print(f"Generating PPTX from {csv_file}...")
    df = pd.read_csv(csv_file)
    
    prs = Presentation()
    
    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Logical Reasoning Robustness"
    subtitle.text = "Experimental Results Summary"
    
    # Results Slide
    slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Accuracy Summary"
    
    rows, cols = df.shape
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(0.8)
    
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
    
    # Set headers
    for i, col_name in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = str(col_name)
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
    # Set data
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r + 1, c)
            val = df.iloc[r, c]
            cell.text = f"{val:.4f}" if isinstance(val, float) else str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(10) # Smaller font for PPT
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    prs.save(output_file)
    print(f"PPTX saved to {output_file}")

if __name__ == "__main__":
    csv_input = "evaluation_summary.csv"
    if os.path.exists(csv_input):
        generate_pdf(csv_input, "robustness_report.pdf")
        generate_pptx(csv_input, "robustness_presentation.pptx")
    else:
        print(f"Error: {csv_input} not found. Run summarize_results.py first.")
